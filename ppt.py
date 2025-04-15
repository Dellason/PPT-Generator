from pptx import Presentation

#Making copy of template
Temp = Presentation('Template.pptx')
Temp.save('yourppt.pptx')
#Selecting the page layout to be duplicated
page = Temp.slides[1]

for shape in page.shapes:
    if shape.has_table:
        table = shape.table
        row_num = 1
        #max_row = len(task_list.people)
        #for person in task_list.people:
        cell =  table.cell(row_num, 0)
        cell.text = "AMA" #person.name

        # Try and create a new slide when the previous is full 
        # Get slide dimensions
        slide_width = Temp.slide_width
        slide_height = Temp.slide_height
        slide_area = slide_width * slide_height

        # Initialize occupied area
        occupied_area = 0

        # Iterate through shapes to calculate occupied area
        for shape in page.shapes:
            shape_area = shape.width * shape.height
            occupied_area += shape_area

        # Define a threshold for "fullness" (e.g., 90% of slide area)
        threshold = 0.9 * slide_area

        # Determine if the slide is full
        if occupied_area >= threshold:
            print("The slide is considered full.")
        else:
            print("The slide has available space.")

Temp.save('yourppt.pptx')      

if shape.has_table:
        current_table = shape.table
        
# for i, row in enumerate(table.rows):
#     row_data = str
#     print()
#             row
#         cell = table.cell(row, )
#         cell.text


# for i, row in enumerate(table.rows):
#     cell = table.cell(row,)
            
        




