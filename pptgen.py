from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
import os
from pydantic import BaseModel
from typing import List
from pptx import Presentation
from copy import deepcopy
from bullets import copy_slide

app = FastAPI(title="PPT Generator")

# How the user input should be for each person
class Task(BaseModel):
    project: str
    name: str
    tasks: List[str]

# How the full user input should look like
class TaskList(BaseModel):
    people: List[Task]  # a list of people and their tasks 

# Actual APP
@app.post("/generate-ppt")
async def generate_ppt(task_list: TaskList):
    if not task_list.people:
        raise HTTPException(status_code=400, detail="No people provided in the input")
    
    # Check if template exists
    if not os.path.exists('Template.pptx'):
        raise HTTPException(status_code=404, detail="Template.pptx not found")
    
    # Create a new presentation
    prs = Presentation('Template.pptx')
    
    # Make sure the presentation has at least 2 slides
    if len(prs.slides) < 2:
        raise HTTPException(status_code=500, detail="Template should have at least 2 slides")
    
    # Get the template slide layout (second slide)
    page_layout = prs.slides[1].slide_layout
    
    # Initialize current slide as the second slide in the template
    current_slide = prs.slides[1]
    
    # Find the table in the template slide
    current_table = None
    for shape in current_slide.shapes:
        if shape.has_table:
            current_table = shape.table
            break
    
    if not current_table:
        raise HTTPException(status_code=500, detail="No table found in the template slide")
    
    # Variables to track table filling
    row_num = 1
    max_row = 5
    compile_task_length = 0
    
    for person in task_list.people:
        tasks_text = "\n".join([f"• {task}" for task in person.tasks])
        compile_task_length += len(tasks_text)
        
        # Check if we need a new slide
        if row_num >= max_row or compile_task_length >= 1200:
            # Create a new slide
            #new_slide = prs.slides.add_slide(page_layout)
            
            # # Copy shapes from the template slide to the new slide
            # for shape in prs.slides[1].shapes:
            #     el = shape.element
            #     newel = deepcopy(el)
            #     new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

            
            # Update current slide reference
            new_slide = copy_slide()
            current_slide += new_slide
            
            # Find the table in the new slide
            current_table = None
            for shape in current_slide.shapes:
                if shape.has_table:
                    current_table = shape.table
                    break
            
            if not current_table:
                raise HTTPException(status_code=500, detail="Failed to find table in new slide")
            
            # Reset counters for the new slide
            row_num = 1
            compile_task_length = 0
        
        # Fill the table cells for this person
        try:
            project_cell = current_table.cell(row_num, 0)
            project_cell.text = person.project
            
            name_cell = current_table.cell(row_num, 1)
            name_cell.text = person.name
            
            progress_cell = current_table.cell(row_num, 2)
            progress_cell.text = tasks_text
            
            # Increment row counter for next person
            row_num += 1
            
        except IndexError:
            raise HTTPException(
                status_code=500, 
                detail=f"Table index out of range. Check table dimensions. Row: {row_num}"
            )
    
    # Save the presentation
    prs.save('tasks_presentation.pptx')
    
    # Return the file as a downloadable response
    return FileResponse(
        path='tasks_presentation.pptx',
        filename="tasks_presentation.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

# Root endpoint for API information
@app.get("/")
async def root():
    return {
        "message": "PPT Generator API is running. Send POST requests to /generate-ppt",
        "example_input": {
            "people": [
                {
                    "name": "Nathan",
                    "project": "AgentX",
                    "tasks": ["Wrote code", "Hosted on render"]
                },
                {
                    "name": "Richmond",
                    "project": "AgentX",
                    "tasks": ["Fixed bugs in flow", "Updated connector"]
                }
            ]
        }
    }
