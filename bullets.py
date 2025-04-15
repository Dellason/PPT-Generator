#To create new slides 
from pptx import Presentation
prs = Presentation('Template.pptx')
working_page = prs.slides[1]

def copy_slide():
    for shape in working_page.shapes:
        el = shape.element
    return working_page
    
 


# def create_new_slide():
#     new_slide = prs.slides.add_slide(template_slide.slide_layout)
    
#     # Copy all shapes 
#     for shape in template_slide.shapes:
#         el = shape.element
#         newel = deepcopy(el)
#         new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
#     # Find and return the table in the new slide
#     for shape in new_slide.shapes:
#         if shape.has_table:
#             return new_slide, shape.table
    
#     return new_slide, None

# # Start with the first data page being the template slide itself
# current_slide = template_slide
# current_table = template_table

# # Track position and content
# row_num = 1
# current_text_length = 0